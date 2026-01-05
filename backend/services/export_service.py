"""Export service for generating PowerPoint and PDF versions of the pitch deck"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from reportlab.lib.pagesizes import letter, landscape
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_CENTER, TA_LEFT
import io
from typing import List, Dict, Any


# NeoNoble brand colors
TEAL_PRIMARY = RGBColor(0x14, 0xB8, 0xA6)  # teal-500
TEAL_DARK = RGBColor(0x0F, 0x76, 0x6E)     # teal-700
SLATE_DARK = RGBColor(0x1E, 0x29, 0x3B)    # slate-800
SLATE_MEDIUM = RGBColor(0x47, 0x55, 0x69)  # slate-600
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def create_pptx(slides: List[Dict[str, Any]], company_info: Dict[str, Any]) -> bytes:
    """Generate a PowerPoint presentation from slides data"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Define layouts
    blank_layout = prs.slide_layouts[6]  # Blank layout
    
    for slide_data in slides:
        slide = prs.slides.add_slide(blank_layout)
        _add_slide_content(slide, slide_data, prs)
    
    # Save to bytes
    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer.getvalue()


def _add_slide_content(slide, slide_data: Dict[str, Any], prs):
    """Add content to a PowerPoint slide based on slide data"""
    # Add header bar
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.8)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = TEAL_DARK
    header_shape.line.fill.background()
    
    # Add company name in header
    company_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(4), Inches(0.4))
    company_tf = company_box.text_frame
    company_p = company_tf.paragraphs[0]
    company_p.text = "NeoNoble Ramp & NeoExchange"
    company_p.font.size = Pt(14)
    company_p.font.bold = True
    company_p.font.color.rgb = WHITE
    
    # Add slide number
    num_box = slide.shapes.add_textbox(Inches(12), Inches(0.2), Inches(1), Inches(0.4))
    num_tf = num_box.text_frame
    num_p = num_tf.paragraphs[0]
    num_p.text = f"{slide_data['id']} / 11"
    num_p.font.size = Pt(12)
    num_p.font.color.rgb = WHITE
    num_p.alignment = PP_ALIGN.RIGHT
    
    # Add subtitle (category)
    subtitle_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(12), Inches(0.4))
    subtitle_tf = subtitle_box.text_frame
    subtitle_p = subtitle_tf.paragraphs[0]
    subtitle_p.text = slide_data['subtitle'].upper()
    subtitle_p.font.size = Pt(12)
    subtitle_p.font.bold = True
    subtitle_p.font.color.rgb = TEAL_PRIMARY
    
    # Add main title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(12), Inches(0.8))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = slide_data['title']
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = SLATE_DARK
    
    # Add content based on slide type
    content = slide_data.get('content', {})
    y_pos = 2.6
    
    # Handle different content types
    if 'headline' in content:
        _add_text_block(slide, content.get('headline', ''), Inches(0.75), Inches(y_pos), Inches(6), Pt(18), SLATE_DARK, bold=True)
        y_pos += 0.6
    
    if 'description' in content:
        _add_text_block(slide, content.get('description', ''), Inches(0.75), Inches(y_pos), Inches(11), Pt(14), SLATE_MEDIUM)
        y_pos += 0.8
    
    if 'vision' in content:
        _add_text_block(slide, content.get('vision', ''), Inches(0.75), Inches(y_pos), Inches(11), Pt(14), SLATE_MEDIUM)
        y_pos += 0.8
    
    if 'keyPoints' in content:
        for point in content['keyPoints']:
            _add_bullet_point(slide, point, Inches(0.75), Inches(y_pos), Inches(11))
            y_pos += 0.45
    
    if 'valueProps' in content:
        for i, prop in enumerate(content['valueProps'][:4]):
            x = Inches(0.75 + (i % 2) * 6)
            y = Inches(y_pos + (i // 2) * 1.2)
            _add_value_prop(slide, prop, x, y)
    
    if 'userSegments' in content:
        for i, segment in enumerate(content['userSegments'][:3]):
            x = Inches(0.75 + i * 4)
            _add_segment_box(slide, segment, x, Inches(y_pos))
    
    if 'useCases' in content:
        y_start = y_pos + 1.5 if 'userSegments' in content else y_pos
        for i, use_case in enumerate(content['useCases']):
            _add_bullet_point(slide, use_case, Inches(0.75), Inches(y_start + i * 0.4), Inches(11))
    
    if 'marketStats' in content:
        for i, stat in enumerate(content['marketStats']):
            x = Inches(7 + (i % 2) * 3)
            y = Inches(y_pos + (i // 2) * 1.2)
            _add_stat_box(slide, stat, x, y)
    
    if 'partnerResponsibilities' in content:
        _add_text_block(slide, "Provider-of-Record Partner:", Inches(0.75), Inches(y_pos), Inches(5), Pt(14), TEAL_DARK, bold=True)
        for i, item in enumerate(content['partnerResponsibilities'][:6]):
            _add_bullet_point(slide, item, Inches(0.75), Inches(y_pos + 0.4 + i * 0.35), Inches(5.5))
    
    if 'neonobleRole' in content:
        _add_text_block(slide, "NeoNoble Role:", Inches(6.75), Inches(y_pos), Inches(5), Pt(14), TEAL_DARK, bold=True)
        for i, item in enumerate(content['neonobleRole'][:5]):
            _add_bullet_point(slide, item, Inches(6.75), Inches(y_pos + 0.4 + i * 0.35), Inches(5.5))
    
    if 'phases' in content:
        for i, phase in enumerate(content['phases'][:3]):
            x = Inches(0.75 + i * 4.1)
            _add_phase_box(slide, phase, x, Inches(y_pos))
    
    if 'callToAction' in content:
        _add_text_block(slide, content.get('callToAction', ''), Inches(0.75), Inches(y_pos), Inches(11), Pt(14), SLATE_MEDIUM)
    
    if 'contact' in content:
        contact = content['contact']
        y_contact = 4.5
        if 'platforms' in contact:
            for i, platform in enumerate(contact['platforms']):
                _add_text_block(slide, f"{platform['name']} — {platform['description']}", Inches(0.75), Inches(y_contact + i * 0.4), Inches(11), Pt(12), SLATE_DARK, bold=True)
        if 'email' in contact:
            _add_text_block(slide, f"Contact: {contact['email']}", Inches(0.75), Inches(y_contact + 1), Inches(11), Pt(14), TEAL_PRIMARY, bold=True)
    
    if 'closing' in content:
        _add_text_block(slide, content.get('closing', ''), Inches(0.75), Inches(6), Inches(11), Pt(12), SLATE_MEDIUM)


def _add_text_block(slide, text: str, left, top, width, font_size, color, bold=False):
    """Add a text block to the slide"""
    box = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold


def _add_bullet_point(slide, text: str, left, top, width):
    """Add a bullet point to the slide"""
    box = slide.shapes.add_textbox(left, top, width, Inches(0.35))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"• {text}"
    p.font.size = Pt(12)
    p.font.color.rgb = SLATE_MEDIUM


def _add_value_prop(slide, prop: Dict[str, str], left, top):
    """Add a value proposition box"""
    _add_text_block(slide, prop.get('title', ''), left, top, Inches(5), Pt(14), SLATE_DARK, bold=True)
    _add_text_block(slide, prop.get('description', ''), left, top + Inches(0.35), Inches(5), Pt(11), SLATE_MEDIUM)


def _add_segment_box(slide, segment: Dict[str, str], left, top):
    """Add a user segment box"""
    _add_text_block(slide, segment.get('segment', ''), left, top, Inches(3.5), Pt(14), SLATE_DARK, bold=True)
    _add_text_block(slide, segment.get('description', ''), left, top + Inches(0.35), Inches(3.5), Pt(11), SLATE_MEDIUM)


def _add_stat_box(slide, stat: Dict[str, str], left, top):
    """Add a market stat box"""
    _add_text_block(slide, stat.get('label', ''), left, top, Inches(2.5), Pt(10), SLATE_MEDIUM)
    _add_text_block(slide, stat.get('value', ''), left, top + Inches(0.25), Inches(2.5), Pt(24), SLATE_DARK, bold=True)
    _add_text_block(slide, stat.get('note', ''), left, top + Inches(0.65), Inches(2.5), Pt(10), TEAL_PRIMARY)


def _add_phase_box(slide, phase: Dict[str, Any], left, top):
    """Add a roadmap phase box"""
    _add_text_block(slide, phase.get('phase', ''), left, top, Inches(3.5), Pt(10), TEAL_PRIMARY, bold=True)
    _add_text_block(slide, phase.get('title', ''), left, top + Inches(0.25), Inches(3.5), Pt(16), SLATE_DARK, bold=True)
    _add_text_block(slide, phase.get('timeline', ''), left, top + Inches(0.55), Inches(3.5), Pt(10), SLATE_MEDIUM)
    for i, item in enumerate(phase.get('items', [])[:4]):
        _add_bullet_point(slide, item, left, top + Inches(0.85 + i * 0.35), Inches(3.5))


def create_pdf(slides: List[Dict[str, Any]], company_info: Dict[str, Any]) -> bytes:
    """Generate a PDF version of the pitch deck"""
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=landscape(letter),
        rightMargin=50,
        leftMargin=50,
        topMargin=40,
        bottomMargin=40
    )
    
    # Custom styles
    styles = getSampleStyleSheet()
    
    styles.add(ParagraphStyle(
        'SlideTitle',
        parent=styles['Heading1'],
        fontSize=28,
        spaceAfter=10,
        textColor=colors.HexColor('#1E293B'),
        fontName='Helvetica-Bold'
    ))
    
    styles.add(ParagraphStyle(
        'SlideSubtitle',
        parent=styles['Normal'],
        fontSize=12,
        spaceAfter=5,
        textColor=colors.HexColor('#14B8A6'),
        fontName='Helvetica-Bold'
    ))
    
    styles.add(ParagraphStyle(
        'SlideBody',
        parent=styles['Normal'],
        fontSize=11,
        spaceAfter=8,
        textColor=colors.HexColor('#475569'),
        leading=16
    ))
    
    styles.add(ParagraphStyle(
        'BulletPoint',
        parent=styles['Normal'],
        fontSize=10,
        spaceAfter=6,
        textColor=colors.HexColor('#475569'),
        leftIndent=20,
        bulletIndent=10
    ))
    
    styles.add(ParagraphStyle(
        'SectionHeader',
        parent=styles['Normal'],
        fontSize=12,
        spaceAfter=8,
        textColor=colors.HexColor('#0F766E'),
        fontName='Helvetica-Bold'
    ))
    
    story = []
    
    for i, slide_data in enumerate(slides):
        # Add slide header
        story.append(Paragraph(f"{slide_data['subtitle'].upper()}", styles['SlideSubtitle']))
        story.append(Paragraph(slide_data['title'], styles['SlideTitle']))
        story.append(Spacer(1, 15))
        
        content = slide_data.get('content', {})
        
        # Handle different content types
        if 'headline' in content:
            story.append(Paragraph(f"<b>{content['headline']}</b>", styles['SlideBody']))
            story.append(Spacer(1, 5))
        
        if 'description' in content:
            story.append(Paragraph(content['description'], styles['SlideBody']))
            story.append(Spacer(1, 10))
        
        if 'vision' in content:
            story.append(Paragraph(content['vision'], styles['SlideBody']))
            story.append(Spacer(1, 10))
        
        if 'keyPoints' in content:
            for point in content['keyPoints']:
                story.append(Paragraph(f"• {point}", styles['BulletPoint']))
        
        if 'valueProps' in content:
            for prop in content['valueProps']:
                story.append(Paragraph(f"<b>{prop['title']}</b>: {prop['description']}", styles['SlideBody']))
        
        if 'userSegments' in content:
            story.append(Paragraph("User Segments:", styles['SectionHeader']))
            for segment in content['userSegments']:
                story.append(Paragraph(f"• <b>{segment['segment']}</b>: {segment['description']}", styles['BulletPoint']))
        
        if 'useCases' in content:
            story.append(Spacer(1, 10))
            story.append(Paragraph("Primary Use Cases:", styles['SectionHeader']))
            for use_case in content['useCases']:
                story.append(Paragraph(f"• {use_case}", styles['BulletPoint']))
        
        if 'primaryMarket' in content:
            story.append(Paragraph(f"<b>Primary Market:</b> {content['primaryMarket']}", styles['SlideBody']))
            story.append(Spacer(1, 10))
        
        if 'marketStats' in content:
            story.append(Paragraph("Market Statistics:", styles['SectionHeader']))
            for stat in content['marketStats']:
                story.append(Paragraph(f"• <b>{stat['label']}</b>: {stat['value']} ({stat['note']})", styles['BulletPoint']))
        
        if 'geographicAdvantages' in content:
            story.append(Spacer(1, 10))
            story.append(Paragraph("Geographic Advantages:", styles['SectionHeader']))
            for advantage in content['geographicAdvantages']:
                story.append(Paragraph(f"• {advantage}", styles['BulletPoint']))
        
        if 'model' in content:
            story.append(Paragraph(f"<b>Model:</b> {content['model']}", styles['SlideBody']))
        
        if 'partnerResponsibilities' in content:
            story.append(Spacer(1, 10))
            story.append(Paragraph("Provider-of-Record Partner Responsibilities:", styles['SectionHeader']))
            for item in content['partnerResponsibilities']:
                story.append(Paragraph(f"• {item}", styles['BulletPoint']))
        
        if 'neonobleRole' in content:
            story.append(Spacer(1, 10))
            story.append(Paragraph("NeoNoble Role:", styles['SectionHeader']))
            for item in content['neonobleRole']:
                story.append(Paragraph(f"• {item}", styles['BulletPoint']))
        
        if 'layers' in content:
            story.append(Paragraph("Technical Layers:", styles['SectionHeader']))
            for layer in content['layers']:
                components = ', '.join(layer['components'])
                story.append(Paragraph(f"• <b>{layer['name']}</b>: {components}", styles['BulletPoint']))
        
        if 'features' in content:
            story.append(Spacer(1, 10))
            story.append(Paragraph("Key Features:", styles['SectionHeader']))
            for feature in content['features']:
                story.append(Paragraph(f"• {feature}", styles['BulletPoint']))
        
        if 'partnershipTypes' in content:
            for ptype in content['partnershipTypes']:
                story.append(Paragraph(f"<b>{ptype['type']}</b>", styles['SectionHeader']))
                story.append(Paragraph(ptype['description'], styles['SlideBody']))
                for benefit in ptype['benefits']:
                    story.append(Paragraph(f"• {benefit}", styles['BulletPoint']))
        
        if 'integrationOptions' in content:
            story.append(Spacer(1, 10))
            story.append(Paragraph("Integration Options:", styles['SectionHeader']))
            for option in content['integrationOptions']:
                story.append(Paragraph(f"• {option}", styles['BulletPoint']))
        
        if 'revenueModel' in content:
            story.append(Paragraph(f"<b>Revenue Model:</b> {content['revenueModel']}", styles['SlideBody']))
        
        if 'projections' in content:
            story.append(Spacer(1, 10))
            story.append(Paragraph("Volume Projections:", styles['SectionHeader']))
            projections = content['projections']
            story.append(Paragraph(f"• Early Stage: {projections.get('earlyStage', '')}", styles['BulletPoint']))
            story.append(Paragraph(f"• Scaling Phase: {projections.get('scalingPhase', '')}", styles['BulletPoint']))
            story.append(Paragraph(f"• Growth Phase: {projections.get('growthPhase', '')}", styles['BulletPoint']))
        
        if 'growthDrivers' in content:
            story.append(Spacer(1, 10))
            story.append(Paragraph("Growth Drivers:", styles['SectionHeader']))
            for driver in content['growthDrivers']:
                story.append(Paragraph(f"• {driver}", styles['BulletPoint']))
        
        if 'phases' in content:
            for phase in content['phases']:
                story.append(Paragraph(f"<b>{phase['phase']}: {phase['title']}</b> ({phase['timeline']})", styles['SectionHeader']))
                for item in phase['items']:
                    story.append(Paragraph(f"• {item}", styles['BulletPoint']))
        
        if 'callToAction' in content:
            story.append(Paragraph(content['callToAction'], styles['SlideBody']))
        
        if 'discussionTopics' in content:
            story.append(Spacer(1, 10))
            story.append(Paragraph("Discussion Topics:", styles['SectionHeader']))
            for topic in content['discussionTopics']:
                story.append(Paragraph(f"• {topic}", styles['BulletPoint']))
        
        if 'contact' in content:
            contact = content['contact']
            story.append(Spacer(1, 15))
            if 'platforms' in contact:
                for platform in contact['platforms']:
                    story.append(Paragraph(f"<b>{platform['name']}</b> — {platform['description']}", styles['SlideBody']))
                    story.append(Paragraph(f"Website: {platform['website']}", styles['BulletPoint']))
            if 'email' in contact:
                story.append(Spacer(1, 10))
                story.append(Paragraph(f"<b>Contact Email:</b> {contact['email']}", styles['SlideBody']))
        
        if 'closing' in content:
            story.append(Spacer(1, 15))
            story.append(Paragraph(f"<i>{content['closing']}</i>", styles['SlideBody']))
        
        # Add page break between slides (except last)
        if i < len(slides) - 1:
            story.append(PageBreak())
    
    doc.build(story)
    buffer.seek(0)
    return buffer.getvalue()
